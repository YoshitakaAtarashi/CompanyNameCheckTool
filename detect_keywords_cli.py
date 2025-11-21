"""
PowerPoint キーワード検出ツール（CLI版）
ディレクトリを指定して、PPTファイル内のキーワードを検出し結果を出力します。
"""

import os
import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation
from datetime import datetime


def load_config():
    """設定ファイルを読み込む"""
    config_file = 'config.json'
    default_config = {
        'default_keywords': ['OldCompany', '旧社名', 'Old Company Name'],
        'allowed_extensions': ['pptx', 'ppt']
    }
    
    if os.path.exists(config_file):
        try:
            with open(config_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"警告: 設定ファイルの読み込みに失敗しました: {str(e)}")
            print("デフォルト設定を使用します。")
            return default_config
    return default_config


def find_ppt_files(directory, recursive=True):
    """指定ディレクトリ内のPPTファイルを検索"""
    ppt_files = []
    path = Path(directory)
    
    if not path.exists():
        print(f"エラー: ディレクトリが存在しません: {directory}")
        return []
    
    if not path.is_dir():
        print(f"エラー: 指定されたパスはディレクトリではありません: {directory}")
        return []
    
    try:
        if recursive:
            # 再帰的に検索
            for ext in ['*.pptx', '*.ppt']:
                ppt_files.extend(path.rglob(ext))
        else:
            # 直下のみ検索
            for ext in ['*.pptx', '*.ppt']:
                ppt_files.extend(path.glob(ext))
        
        # 隠しファイルを除外
        ppt_files = [f for f in ppt_files if not any(part.startswith('.') for part in f.parts)]
        
    except Exception as e:
        print(f"エラー: ファイル検索中にエラーが発生しました: {str(e)}")
        return []
    
    return sorted(ppt_files)


def find_keywords_in_presentation(prs, keywords):
    """プレゼンテーション内のキーワードを検出
    通常スライドとマスタースライドの両方をチェック"""
    results = []
    
    # 通常スライドを処理
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape_num, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                found_keywords = []
                total_count = 0
                
                # すべてのキーワードを検査
                for keyword in keywords:
                    if keyword.lower() in shape.text.lower():
                        count = shape.text.lower().count(keyword.lower())
                        found_keywords.append(keyword)
                        total_count += count
                
                # いずれかのキーワードが見つかった場合
                if found_keywords:
                    results.append({
                        'slide': slide_num,
                        'shape': shape_num,
                        'text': shape.text[:100],  # 最初の100文字のみ
                        'keywords': found_keywords,
                        'count': total_count,
                        'is_master': False
                    })
    
    # マスタースライドを処理（複数のマスターグループに対応）
    try:
        for master_group_num, slide_master in enumerate(prs.slide_masters):
            for layout_num, layout in enumerate(slide_master.slide_layouts):
                for shape_num, shape in enumerate(layout.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        found_keywords = []
                        total_count = 0
                        
                        for keyword in keywords:
                            if keyword.lower() in shape.text.lower():
                                count = shape.text.lower().count(keyword.lower())
                                found_keywords.append(keyword)
                                total_count += count
                        
                        if found_keywords:
                            results.append({
                                'slide': f'Master Group {master_group_num + 1}, Layout {layout_num + 1}',
                                'shape': shape_num,
                                'text': shape.text[:100],  # 最初の100文字のみ
                                'keywords': found_keywords,
                                'count': total_count,
                                'is_master': True
                            })
    except Exception as e:
        print(f"    警告: マスタースライド処理エラー: {str(e)}")
    
    return results


def detect_keywords_in_file(file_path, keywords):
    """1つのファイル内のキーワードを検出"""
    try:
        prs = Presentation(str(file_path))
        results = find_keywords_in_presentation(prs, keywords)
        return {
            'success': True,
            'results': results,
            'error': None
        }
    except Exception as e:
        return {
            'success': False,
            'results': [],
            'error': str(e)
        }


def format_results_text(all_results, target_directory, show_all_files=False):
    """検出結果をシンプルなリスト形式で整形"""
    output = []
    
    total_files = len(all_results)
    files_with_keywords = sum(1 for r in all_results if r['results'])
    
    # ファイルパスと検出数のリスト
    for file_result in all_results:
        file_path = file_result['file']
        
        if file_result['success']:
            detection_count = len(file_result['results'])
            # show_all_filesがTrueの場合は全ファイル、Falseの場合は検出があったファイルのみ
            if show_all_files or detection_count > 0:
                output.append(f"{file_path}\t{detection_count}")
        else:
            # エラーの場合は常に表示
            output.append(f"{file_path}\t0\t(エラー: {file_result['error']})")
    
    # サマリー情報
    output.append("")
    output.append("=" * 80)
    output.append(f"対象ディレクトリ: {target_directory}")
    output.append(f"検出ファイル数: {files_with_keywords}/{total_files}")
    output.append(f"実施日時: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    output.append("=" * 80)
    
    return "\n".join(output)


def save_results_to_file(content, output_file):
    """結果をファイルに保存"""
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(content)
        print(f"\n結果を保存しました: {output_file}")
    except Exception as e:
        print(f"\nエラー: ファイル保存に失敗しました: {str(e)}")


def main():
    parser = argparse.ArgumentParser(
        description='PowerPointファイル内のキーワードを検出します',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
使用例:
  python detect_keywords_cli.py C:\\Documents\\Presentations
  python detect_keywords_cli.py C:\\Documents --no-recursive
  python detect_keywords_cli.py C:\\Documents --keywords "OldCompany" "旧社名"
  python detect_keywords_cli.py C:\\Documents --output results.txt
        """
    )
    
    parser.add_argument('directory', help='検索対象のディレクトリパス')
    parser.add_argument('--keywords', '-k', nargs='+', help='検索するキーワード（スペース区切り）')
    parser.add_argument('--no-recursive', '-n', action='store_true', 
                       help='サブディレクトリを検索しない')
    parser.add_argument('--output', '-o', help='結果を保存するファイル名')
    parser.add_argument('--show-all', '-a', action='store_true',
                       help='検出数が0のファイルも含めて全ファイルを表示')
    
    args = parser.parse_args()
    
    # 設定読み込み
    config = load_config()
    
    # キーワード設定
    keywords = args.keywords if args.keywords else config['default_keywords']
    
    print("=" * 80)
    print("PowerPoint キーワード検出ツール (CLI版)")
    print("=" * 80)
    print(f"検索ディレクトリ: {args.directory}")
    print(f"検索キーワード: {', '.join(keywords)}")
    print(f"再帰検索: {'いいえ' if args.no_recursive else 'はい'}")
    print("-" * 80)
    
    # PPTファイルを検索
    print("\nPPTファイルを検索中...")
    ppt_files = find_ppt_files(args.directory, recursive=not args.no_recursive)
    
    if not ppt_files:
        print("PPTファイルが見つかりませんでした。")
        return
    
    print(f"{len(ppt_files)} 件のPPTファイルが見つかりました。\n")
    
    # 各ファイルを処理
    all_results = []
    for i, file_path in enumerate(ppt_files, 1):
        print(f"[{i}/{len(ppt_files)}] 検査中: {file_path.name} ... ", end='', flush=True)
        
        result = detect_keywords_in_file(file_path, keywords)
        result['file'] = str(file_path)
        all_results.append(result)
        
        if result['success']:
            if result['results']:
                print(f"✓ {len(result['results'])} 箇所で検出")
            else:
                print("検出なし")
        else:
            print(f"✗ エラー")
    
    # 結果を整形
    output_text = format_results_text(all_results, args.directory, show_all_files=args.show_all)
    
    # 結果を表示
    print("\n")
    print(output_text)
    
    # ファイルに保存
    if args.output:
        save_results_to_file(output_text, args.output)
    
    # エラーがあった場合は終了コード1
    if any(not r['success'] for r in all_results):
        sys.exit(1)


if __name__ == '__main__':
    try:
        main()
    except KeyboardInterrupt:
        print("\n\n処理を中断しました。")
        sys.exit(130)
    except Exception as e:
        print(f"\n予期しないエラーが発生しました: {str(e)}")
        sys.exit(1)
